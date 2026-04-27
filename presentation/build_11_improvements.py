"""Topic 11 — Further improvements possible.

Built from docs/TECHNICAL_GUIDE.md §11 (unwired IaC) and §13 (known limitations),
plus open user stories US33-35 from the agile backlog.
"""
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import _pptx_helpers as h


def RGBColor_safe(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def improvements(prs):
    s = h.blank(prs)
    h.header_band(s, "11.  Further improvements possible",
                  "From technical debt to product roadmap")

    h.add_text(s,
               "We documented every known limitation in TECHNICAL_GUIDE.md §13 — and the repo "
               "already carries an unwired Bicep + GitHub Actions path so a fresh tenant rebuild "
               "is a few approvals away, not a re-write.",
               left=Inches(0.5), top=Inches(1.25), width=Inches(12.3), height=Inches(0.55),
               size=12, color=h.GREY)

    # 4 themed columns
    columns = [
        ("Reproducibility & DR", h.NAVY, h.LIGHT_BLUE, [
            "Activate deploy-sql.yml (needs Key Vault Secrets User on UAMI).",
            "Wire infrastructure/future/workflows/ → CI for full Bicep rebuild.",
            "Bacpac seed for fresh DBs already in infrastructure/exported/.",
        ]),
        ("Code & data quality", h.RED, RGBColor_safe(0xFF, 0xEB, 0xEE), [
            "Unify the tariff source — currently triplicated (SQL computed col, JSON, ref_electricity_tariff SCD2).",
            "Add a LoadBatchId / IngestedAt to fact tables — currently no lineage.",
            "Wire PL_Bronze_MeteoFuture into PL_Ingest_Bronze (today it's orphaned).",
            "Make notebook paths in ADF activities user-agnostic — they hard-code /Repos/<user>/.",
        ]),
        ("Security & GDPR", h.SERVE_GREEN, h.SERVE_FILL, [
            "Finish RLS rollout (US#26 / US#27 still in progress).",
            "Provide a written GDPR compliance assessment (US#35).",
            "Populate ref_user_division_access for all Directors.",
        ]),
        ("Product roadmap", h.GOLD, h.GOLD_FILL, [
            "Internationalisation of dashboards & stories (US#33).",
            "Scalability assessment + capacity forecast (US#34).",
            "Re-train ML models on a longer window (current data ends 2023-04-19).",
            "Promote vw_prediction_accuracy as a first-class KPI surfaced in BI.",
        ]),
    ]

    col_w = Inches(3.05)
    gap = Inches(0.15)
    top = Inches(2.00)
    left = Inches(0.30)
    for title, accent, fill, items in columns:
        h.rounded_card(s, left=left, top=top, width=col_w, height=Inches(0.55),
                       fill=accent, line=accent)
        h.add_text(s, title, left=left, top=top, width=col_w, height=Inches(0.55),
                   size=14, bold=True, color=h.WHITE,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        h.rounded_card(s, left=left, top=top + Inches(0.65),
                       width=col_w, height=Inches(4.45),
                       fill=fill, line=accent)
        ry = top + Inches(0.80)
        for it in items:
            h.add_text(s, "•  " + it,
                       left=left + Inches(0.18), top=ry,
                       width=col_w - Inches(0.36), height=Inches(1.10),
                       size=11, color=h.DARK)
            ry += Inches(1.05)
        left += col_w + gap

    h.footer(s)


def closing(prs):
    s = h.blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, h.SLIDE_W, h.SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = h.NAVY; bg.line.fill.background()

    ribbon = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4), h.SLIDE_W, Inches(0.10))
    ribbon.fill.solid(); ribbon.fill.fore_color.rgb = h.GOLD; ribbon.line.fill.background()

    h.add_text(s, "Thank you",
               left=Inches(0.7), top=Inches(2.4), width=Inches(12), height=Inches(1.0),
               size=60, bold=True, color=h.WHITE, align=PP_ALIGN.LEFT)
    h.add_text(s, "Questions?",
               left=Inches(0.7), top=Inches(3.7), width=Inches(12), height=Inches(0.6),
               size=28, color=h.GOLD, align=PP_ALIGN.LEFT)

    h.add_text(s,
               "Repository  ·  github.com/sandersdHES/ADF_DataCycleProject\n"
               "Architecture diagram  ·  docs/ARCHITECTURE.md\n"
               "Technical guide  ·  docs/TECHNICAL_GUIDE.md",
               left=Inches(0.7), top=Inches(5.0), width=Inches(12), height=Inches(1.5),
               size=14, color=h.WHITE, align=PP_ALIGN.LEFT)


if __name__ == "__main__":
    prs = h.open_or_create()
    improvements(prs)
    closing(prs)
    h.save(prs)
    print("ok — improvements + closing slides appended")
