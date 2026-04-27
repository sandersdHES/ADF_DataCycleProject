"""Helpers for building the Bellevue Data Cycle deck.

Kept intentionally small: a few palette constants, geometry helpers, and slide
builders that the per-section scripts call. Run order:

    1. build_01_intro.py        -> writes Bellevue_DataCycle.pptx (cover + intro)
    2. build_02_team.py         -> appends team slide
    3. build_03_organisation.py
    4. build_04_architecture.py
    5. build_05_dataflow.py
    6. build_06_silver.py
    7. build_07_gold.py
    8. build_08_ml.py
    9. build_09_powerbi.py
   10. build_10_sac.py
   11. build_11_improvements.py

Each step is idempotent in the sense that it opens the existing file, appends
its slides, and saves again — so re-running a step duplicates its slides.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

DECK_PATH = Path(__file__).parent / "Bellevue_DataCycle.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette — inspired by the Mermaid medallion colours used in docs/ARCHITECTURE.md
NAVY = RGBColor(0x0D, 0x47, 0xA1)
BLUE = RGBColor(0x19, 0x76, 0xD2)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
BRONZE = RGBColor(0xE6, 0x51, 0x00)
BRONZE_FILL = RGBColor(0xFF, 0xE0, 0xB2)
SILVER = RGBColor(0x45, 0x5A, 0x64)
SILVER_FILL = RGBColor(0xEC, 0xEF, 0xF1)
GOLD = RGBColor(0xFF, 0xA0, 0x00)
GOLD_FILL = RGBColor(0xFF, 0xF8, 0xE1)
ML_PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
ML_FILL = RGBColor(0xF3, 0xE5, 0xF5)
SERVE_GREEN = RGBColor(0x2E, 0x7D, 0x32)
SERVE_FILL = RGBColor(0xE8, 0xF5, 0xE9)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)
DARK = RGBColor(0x21, 0x21, 0x21)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def open_or_create() -> Presentation:
    if DECK_PATH.exists():
        return Presentation(str(DECK_PATH))
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def save(prs: Presentation) -> None:
    prs.save(str(DECK_PATH))


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------- text helpers ----------

def _set_run(run, *, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    text,
    *,
    left,
    top,
    width,
    height,
    size=18,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_bullets(slide, items, *, left, top, width, height, size=16, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"•  {item}"
        _set_run(run, size=size, color=color)
    return box


def header_band(slide, title, subtitle=None, color=NAVY):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.0))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = color
    add_text(
        slide,
        title,
        left=Inches(0.5),
        top=Inches(0.10),
        width=Inches(12.5),
        height=Inches(0.55),
        size=28,
        bold=True,
        color=WHITE,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            left=Inches(0.5),
            top=Inches(0.58),
            width=Inches(12.5),
            height=Inches(0.40),
            size=14,
            color=WHITE,
        )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(1.0), SLIDE_W, Inches(0.05)
    )
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD


def footer(slide, text="Bellevue Data Cycle  •  HES-SO Bellevue Campus  •  ADF_DataCycleProject"):
    add_text(
        slide,
        text,
        left=Inches(0.3),
        top=Inches(7.10),
        width=Inches(12.7),
        height=Inches(0.3),
        size=10,
        color=GREY,
    )


def rounded_card(slide, *, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    return shape


def card(slide, *, left, top, width, height, title, body, fill=LIGHT_BLUE, accent=BLUE):
    rounded_card(slide, left=left, top=top, width=width, height=height, fill=fill, line=accent)
    add_text(
        slide,
        title,
        left=left + Inches(0.15),
        top=top + Inches(0.10),
        width=width - Inches(0.30),
        height=Inches(0.45),
        size=15,
        bold=True,
        color=accent,
    )
    add_text(
        slide,
        body,
        left=left + Inches(0.15),
        top=top + Inches(0.55),
        width=width - Inches(0.30),
        height=height - Inches(0.65),
        size=12,
        color=DARK,
    )


def arrow(slide, *, left, top, width, height=Inches(0.35), color=GREY):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


def chip(slide, *, left, top, width, height, text, fill, color=DARK, size=12, bold=True):
    rounded_card(slide, left=left, top=top, width=width, height=height, fill=fill)
    add_text(
        slide,
        text,
        left=left,
        top=top,
        width=width,
        height=height,
        size=size,
        bold=bold,
        color=color,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
