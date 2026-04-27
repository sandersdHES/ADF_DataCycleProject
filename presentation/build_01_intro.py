"""Cover slide + Topic 1 — Introduction."""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import _pptx_helpers as h


def cover(prs):
    s = h.blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, h.SLIDE_W, h.SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = h.NAVY
    bg.line.fill.background()

    # gold ribbon
    ribbon = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4), h.SLIDE_W, Inches(0.10))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = h.GOLD
    ribbon.line.fill.background()

    h.add_text(s, "BELLEVUE DATA CYCLE",
               left=Inches(0.7), top=Inches(1.6), width=Inches(12), height=Inches(0.7),
               size=14, bold=True, color=h.GOLD, align=PP_ALIGN.LEFT)
    h.add_text(s, "An end-to-end Azure data platform",
               left=Inches(0.7), top=Inches(2.1), width=Inches(12), height=Inches(0.9),
               size=44, bold=True, color=h.WHITE, align=PP_ALIGN.LEFT)
    h.add_text(s, "Ingestion  ·  Medallion Lakehouse  ·  ML Forecasting  ·  BI",
               left=Inches(0.7), top=Inches(2.95), width=Inches(12), height=Inches(0.5),
               size=20, color=h.WHITE, align=PP_ALIGN.LEFT)

    h.add_text(s, "HES-SO Bellevue Campus  ·  Energy Monitoring Project",
               left=Inches(0.7), top=Inches(3.7), width=Inches(12), height=Inches(0.5),
               size=18, bold=True, color=h.WHITE, align=PP_ALIGN.LEFT)

    # tech chips
    chips = ["Azure Data Factory", "Databricks (PySpark)", "ADLS Gen2", "Azure SQL",
             "KNIME Server (GBT)", "Power BI", "SAP Analytics Cloud"]
    chip_w, chip_h, gap = Inches(1.66), Inches(0.45), Inches(0.10)
    total_w = chip_w * len(chips) + gap * (len(chips) - 1)
    start_left = (h.SLIDE_W - total_w) / 2
    for i, label in enumerate(chips):
        h.chip(s, left=start_left + (chip_w + gap) * i, top=Inches(5.4),
               width=chip_w, height=chip_h,
               text=label, fill=h.WHITE, color=h.NAVY, size=11, bold=True)

    h.add_text(s, "Project Presentation  ·  Sprint 6",
               left=Inches(0.7), top=Inches(6.6), width=Inches(12), height=Inches(0.4),
               size=14, color=h.GOLD, align=PP_ALIGN.LEFT, bold=True)


def agenda(prs):
    s = h.blank(prs)
    h.header_band(s, "Agenda", "What we will cover today")

    items_left = [
        ("1", "Introduction", "Project context and objectives"),
        ("2", "Team Presentation", "Who we are"),
        ("3", "How we organised the work", "Agile cadence, sprints, tooling"),
        ("4", "Architecture", "End-to-end system view"),
        ("5", "Data Flow", "Daily orchestration & triggers"),
        ("6", "Silver Transformation", "Bronze → Silver cleansing"),
    ]
    items_right = [
        ("7", "Gold Transformation", "Silver → Gold dims & facts"),
        ("8", "Machine Learning", "KNIME GBT forecasts"),
        ("9", "Power BI Dashboards", "Operational BI"),
        ("10", "SAC Dashboards", "SAP Analytics Cloud"),
        ("11", "Further improvements", "Roadmap"),
    ]

    def render(col_items, x):
        y = Inches(1.4)
        for n, title, sub in col_items:
            badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.7), Inches(0.7))
            badge.fill.solid(); badge.fill.fore_color.rgb = h.NAVY; badge.line.fill.background()
            h.add_text(s, n, left=x, top=y, width=Inches(0.7), height=Inches(0.7),
                       size=20, bold=True, color=h.WHITE,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            h.add_text(s, title, left=x + Inches(0.85), top=y - Inches(0.02),
                       width=Inches(5.2), height=Inches(0.4),
                       size=18, bold=True, color=h.DARK)
            h.add_text(s, sub, left=x + Inches(0.85), top=y + Inches(0.32),
                       width=Inches(5.2), height=Inches(0.4),
                       size=12, color=h.GREY)
            y += Inches(0.85)

    render(items_left, Inches(0.6))
    render(items_right, Inches(7.0))
    h.footer(s)


def introduction(prs):
    s = h.blank(prs)
    h.header_band(s, "1.  Introduction", "Bellevue Data Cycle — context & objectives")

    # left: context paragraph
    h.add_text(s, "Project context",
               left=Inches(0.5), top=Inches(1.3), width=Inches(6.3), height=Inches(0.4),
               size=18, bold=True, color=h.NAVY)
    h.add_bullets(s, [
        "HES-SO Bellevue campus runs solar PV, energy meters, weather sensors and a room-booking system.",
        "Data lives on-premises (SMB / SFTP shares on a Windows VM, IP 10.130.25.152).",
        "Goal — turn this raw operational data into reliable analytics and forecasts that staff and management can actually use.",
    ], left=Inches(0.5), top=Inches(1.8), width=Inches(6.3), height=Inches(2.3), size=14)

    h.add_text(s, "Why a Medallion lakehouse?",
               left=Inches(0.5), top=Inches(4.1), width=Inches(6.3), height=Inches(0.4),
               size=18, bold=True, color=h.NAVY)
    h.add_bullets(s, [
        "Bronze keeps raw fidelity & traceability of source files.",
        "Silver isolates cleaning, GDPR masking and dedup logic in one place.",
        "Gold is a star schema in Azure SQL — ready for BI, ML and SAC consumption.",
    ], left=Inches(0.5), top=Inches(4.5), width=Inches(6.3), height=Inches(2.4), size=14)

    # right: objectives card stack
    h.add_text(s, "Deliverables",
               left=Inches(7.2), top=Inches(1.3), width=Inches(5.8), height=Inches(0.4),
               size=18, bold=True, color=h.NAVY)

    cards = [
        ("Daily ingestion", "4 sources ingested incrementally via ADF + Self-Hosted IR.", h.LIGHT_BLUE, h.BLUE),
        ("Lakehouse + DWH", "Bronze/Silver in ADLS Gen2, Gold in Azure SQL (8 dims, 7 facts, 7 views).", h.GOLD_FILL, h.GOLD),
        ("ML forecasts", "Solar production & building consumption — KNIME Gradient Boosted Trees.", h.ML_FILL, h.ML_PURPLE),
        ("BI delivery", "Power BI reports + SAC dashboards, with RBAC & RLS.", h.SERVE_FILL, h.SERVE_GREEN),
    ]
    top = Inches(1.8)
    for title, body, fill, accent in cards:
        h.card(s, left=Inches(7.2), top=top, width=Inches(5.8), height=Inches(1.15),
               title=title, body=body, fill=fill, accent=accent)
        top += Inches(1.25)

    h.footer(s)


if __name__ == "__main__":
    prs = h.open_or_create()
    cover(prs)
    agenda(prs)
    introduction(prs)
    h.save(prs)
    print("ok — intro pack added", h.DECK_PATH)
