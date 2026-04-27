"""Topic 3 — How we organised the work.

Built from docs/agile_backlog_corrected.md (sprint summary) and the GitHub
Actions / branch model that the repo already uses.
"""
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import _pptx_helpers as h

SPRINTS = [
    ("Sprint 1", "17 Feb → 02 Mar", "Architecture defined, environment ready",  8, "Done"),
    ("Sprint 2", "03 Mar → 16 Mar", "Core ingestion pipelines (solar, conso, weather)", 24, "Done"),
    ("Sprint 3", "17 Mar → 30 Mar", "Remaining ingestion, security, monitoring", 37, "Done"),
    ("Sprint 4", "31 Mar → 13 Apr", "Silver layer + Gold (OLAP) DB built", 34, "Done"),
    ("Sprint 5", "14 Apr → 27 Apr", "Power BI & SAC dashboards + RLS", 42, "In progress"),
    ("Sprint 6", "28 Apr → 08 May", "ML, predictions, docs, deployment, presentation", 83, "Planned"),
]


def slide(prs):
    s = h.blank(prs)
    h.header_band(s, "3.  How we organised the work",
                  "Agile Scrum  ·  6 sprints of 2 weeks  ·  39 user stories  ·  ~228 SP")

    # left — practices column
    h.add_text(s, "Practices",
               left=Inches(0.5), top=Inches(1.25), width=Inches(4.5), height=Inches(0.4),
               size=18, bold=True, color=h.NAVY)
    h.add_bullets(s, [
        "MoSCoW prioritisation per user story.",
        "Two-week sprints, sprint reviews & retros.",
        "User stories sized in story points (planning poker).",
        "Definition of Done shared across the team.",
        "Daily standups during build sprints.",
    ], left=Inches(0.5), top=Inches(1.7), width=Inches(4.5), height=Inches(2.5), size=13)

    h.add_text(s, "Tooling",
               left=Inches(0.5), top=Inches(4.30), width=Inches(4.5), height=Inches(0.4),
               size=18, bold=True, color=h.NAVY)
    h.add_bullets(s, [
        "GitHub repo + Issues / Projects board.",
        "Pull Requests with review (validate.yml gate).",
        "GitHub Actions (deploy-adf.yml) on push to main — OIDC, no secrets.",
        "Excel agile plan exported to Markdown / JSON for traceability.",
        "Architecture diagrams in Mermaid (docs/ARCHITECTURE.md).",
    ], left=Inches(0.5), top=Inches(4.75), width=Inches(4.5), height=Inches(2.3), size=13)

    # right — sprint timeline
    h.add_text(s, "Sprint roadmap",
               left=Inches(5.4), top=Inches(1.25), width=Inches(7.5), height=Inches(0.4),
               size=18, bold=True, color=h.NAVY)

    row_top = Inches(1.75)
    row_h = Inches(0.78)
    label_w = Inches(1.2)
    bar_left = Inches(6.85)
    bar_max = Inches(5.6)
    max_sp = max(sp for _, _, _, sp, _ in SPRINTS)

    status_color = {"Done": h.SERVE_GREEN, "In progress": h.GOLD, "Planned": h.GREY}

    for i, (name, period, goal, sp, status) in enumerate(SPRINTS):
        y = row_top + row_h * i
        # left badge
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(5.4), y, label_w, Inches(0.55))
        badge.fill.solid(); badge.fill.fore_color.rgb = h.NAVY; badge.line.fill.background()
        h.add_text(s, name, left=Inches(5.4), top=y, width=label_w, height=Inches(0.55),
                   size=12, bold=True, color=h.WHITE,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # period
        h.add_text(s, period,
                   left=Inches(5.4), top=y + Inches(0.55), width=label_w, height=Inches(0.22),
                   size=9, color=h.GREY, align=PP_ALIGN.CENTER)
        # bar
        bar_w = int(bar_max * (sp / max_sp))
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 bar_left, y + Inches(0.05),
                                 bar_w, Inches(0.45))
        bar.fill.solid(); bar.fill.fore_color.rgb = status_color[status]
        bar.line.fill.background()
        h.add_text(s, f"  {goal}",
                   left=bar_left, top=y + Inches(0.05),
                   width=bar_w, height=Inches(0.45),
                   size=10, bold=True, color=h.WHITE,
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        # sp tag
        h.add_text(s, f"{sp} SP  ·  {status}",
                   left=bar_left, top=y + Inches(0.50), width=bar_max, height=Inches(0.22),
                   size=9, color=h.GREY)

    # legend
    legend_y = row_top + row_h * len(SPRINTS) + Inches(0.10)
    cx = Inches(6.85)
    for label, col in [("Done", h.SERVE_GREEN), ("In progress", h.GOLD), ("Planned", h.GREY)]:
        sw = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, legend_y, Inches(0.25), Inches(0.18))
        sw.fill.solid(); sw.fill.fore_color.rgb = col; sw.line.fill.background()
        h.add_text(s, label, left=cx + Inches(0.30), top=legend_y - Inches(0.02),
                   width=Inches(1.2), height=Inches(0.25),
                   size=10, color=h.DARK)
        cx += Inches(1.6)

    h.footer(s)


if __name__ == "__main__":
    prs = h.open_or_create()
    slide(prs)
    h.save(prs)
    print("ok — organisation slide appended")
